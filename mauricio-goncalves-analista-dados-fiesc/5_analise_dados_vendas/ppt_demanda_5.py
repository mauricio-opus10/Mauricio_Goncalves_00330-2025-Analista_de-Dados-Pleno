import nbformat
import nbconvert
import os
import re
import base64
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image
import io
import pandas as pd


def extract_graphic_title(cell_source, default_title="Análise de Dados"):
    """
    Tenta extrair um título significativo para o gráfico baseado no código ou comentários
    """
    # Padrões para extrair títulos de gráficos
    title_patterns = [
        r'plt\.title\([\'"](.+?)[\'"]\)',  # Título direto do matplotlib
        r'# *Título *[:]*\s*(.+)',  # Comentário de título
        r'# *Análise *[:]*\s*(.+)',  # Comentário de análise
        r'# *(.+gráfico.+)',  # Comentário contendo a palavra "gráfico"
    ]

    for pattern in title_patterns:
        match = re.search(pattern, cell_source, re.IGNORECASE)
        if match:
            return match.group(1).strip()

    return default_title


def convert_notebook_to_professional_ppt(notebook_path, output_path=None):
    """
    Converte um notebook Jupyter (.ipynb) para uma apresentação PowerPoint profissional 
    focada em resultados de análise de dados
    """
    # Verificar se o arquivo existe
    if not os.path.exists(notebook_path):
        raise FileNotFoundError(
            f"O arquivo {notebook_path} não foi encontrado.")

    # Ler o notebook
    try:
        with open(notebook_path, 'r', encoding='utf-8') as f:
            nb = nbformat.read(f, as_version=4)
    except Exception as e:
        print(f"Erro ao ler o notebook: {e}")
        return

    # Criar uma nova apresentação
    prs = Presentation()

    # Definir layouts de slide
    title_slide_layout = prs.slide_layouts[0]
    section_slide_layout = prs.slide_layouts[1]
    content_slide_layout = prs.slide_layouts[6]

    # Slide de Título
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Análise de Vendas FIESC"
    subtitle.text = "Relatório Executivo para a Área de Marketing"

    # Processar células do notebook para identificar contextos importantes
    visualization_data = []

    for cell in nb.cells:
        # Processar outputs de imagens e gráficos
        if hasattr(cell, 'outputs'):
            for output in cell.outputs:
                if output.output_type in ['display_data', 'execute_result']:
                    # Capturar imagens PNG
                    if 'image/png' in output.data:
                        try:
                            # Decodificar imagem
                            img_data = base64.b64decode(
                                output.data['image/png'])
                            img = Image.open(io.BytesIO(img_data))

                            # Gerar nome de arquivo único
                            img_filename = f'plot_{len(visualization_data)}.png'
                            img_path = os.path.join(
                                os.path.dirname(notebook_path), img_filename)
                            img.save(img_path)

                            # Extrair título do gráfico
                            graph_title = extract_graphic_title(cell.source)

                            visualization_data.append({
                                'path': img_path,
                                'title': graph_title
                            })
                        except Exception as e:
                            print(f"Erro ao processar imagem: {e}")

    # Slide de Preparação dos Dados
    slide = prs.slides.add_slide(section_slide_layout)
    if slide.shapes.title:
        slide.shapes.title.text = "Preparação dos Dados"

    body = slide.shapes.placeholders[1]
    text_frame = body.text_frame
    text_frame.text = (
        "Etapas de Tratamento:\n"
        "- Inspeção da estrutura de dados\n"
        "- Identificação e correção de anomalias\n"
        "- Transformação de tipos de dados\n"
        "- Tratamento de valores ausentes"
    )

    # Adicionar visualizações com títulos
    for viz in visualization_data:
        # Adicionar slide com imagem
        slide = prs.slides.add_slide(content_slide_layout)

        # Adicionar título de forma mais robusta
        # Tentar adicionar título ao placeholder de título
        title_placeholders = [
            shape for shape in slide.shapes if shape.has_text_frame and shape.is_placeholder]
        title_placeholders = [
            ph for ph in title_placeholders if ph.placeholder_format.type == 1]  # Tipo 1 é título

        if title_placeholders:
            title_placeholders[0].text_frame.text = viz['title']
        else:
            # Se não houver placeholder de título, adicionar texto como uma caixa de texto
            left = top = Inches(0.5)
            width = Inches(9)
            height = Inches(1)
            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.text = viz['title']
            text_frame.paragraphs[0].font.size = Pt(18)
            text_frame.paragraphs[0].font.bold = True

        # Adicionar imagem
        left = top = Inches(1)
        pic = slide.shapes.add_picture(viz['path'], left, top, width=Inches(6))

    # Slide de Conclusões
    slide = prs.slides.add_slide(section_slide_layout)
    if slide.shapes.title:
        slide.shapes.title.text = "Conclusões e Recomendações"

    body = slide.shapes.placeholders[1]
    text_frame = body.text_frame
    text_frame.text = (
        "Principais Insights:\n"
        "- [Insight 1]\n"
        "- [Insight 2]\n"
        "- [Insight 3]\n\n"
        "Recomendações para Marketing:\n"
        "- Estratégia de segmentação\n"
        "- Potenciais ações de vendas\n"
        "- Proposta de valor ajustada"
    )

    # Definir caminho de saída
    if output_path is None:
        output_path = os.path.splitext(notebook_path)[0] + '_apresentacao.pptx'

    # Salvar apresentação
    try:
        prs.save(output_path)
        print(f"Apresentação profissional salva em: {output_path}")

        # Limpar imagens temporárias
        for viz in visualization_data:
            try:
                os.remove(viz['path'])
            except Exception as e:
                print(f"Erro ao remover imagem temporária: {e}")

    except Exception as e:
        print(f"Erro ao salvar a apresentação: {e}")


# Caminho do notebook
notebook_path = r"C:\Users\gonca\OneDrive\009 PERSONAL\CV\Processos Seletivos\mauricio-goncalves-analista-dados-fiesc\5_analise_dados_vendas\notebook_analise de vendas_5.ipynb"

# Chamar a função
convert_notebook_to_professional_ppt(notebook_path)
